#!/usr/bin/env python3

####################################################
# ERB Generator Python Script v3.3
# Designed and Written by Salvador Melendez
# GUI Created using: PyQt5 UI code generator 5.14.2
# WARNING! Any changes made to this file can
#          damage the functionality of the script
####################################################


import os
import glob
import xml.etree.ElementTree as ET
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PyQt5 import QtCore, QtGui, QtWidgets
from PyQt5.QtWidgets import QApplication, QWidget, QLabel, QMessageBox, QFileDialog
from PyQt5.QtGui import QIcon, QPixmap
from PyQt5.QtCore import QDate, QTime, QDateTime, Qt


#VARIABLES
order_list = []
findings_file = 'findings.txt'
desktop_dir = os.path.expanduser("~/Desktop")
cwd = os.getcwd()
current_erb = cwd + '/' + 'erb/findings.xml'
data_folder = ''
data_source = ''
findings = {}
event = []
folder_list = ''
f_folder = ''


class Ui_MainWindow(object):
    global order_list, data_folder, findings, event, folder_list
    def setupUi(self, MainWindow):
        MainWindow.setObjectName("MainWindow")
        MainWindow.setEnabled(True)
        MainWindow.resize(1345, 791)
        MainWindow.setMinimumSize(QtCore.QSize(1345, 791))
        MainWindow.setMaximumSize(QtCore.QSize(1345, 791))
        self.centralwidget = QtWidgets.QWidget(MainWindow)
        self.centralwidget.setObjectName("centralwidget")
        self.stackedWidget0 = QtWidgets.QStackedWidget(self.centralwidget)
        self.stackedWidget0.setGeometry(QtCore.QRect(0, 0, 1331, 751))
        self.stackedWidget0.setObjectName("stackedWidget0")
        self.page1 = QtWidgets.QWidget()
        self.page1.setObjectName("page1")
        self.startDate = QtWidgets.QDateEdit(self.page1)
        self.startDate.setGeometry(QtCore.QRect(233, 611, 111, 31))
        self.startDate.setDate(QtCore.QDate(2020, 1, 1))
        self.startDate.setObjectName("startDate")
        self.leadName_label = QtWidgets.QLabel(self.page1)
        self.leadName_label.setGeometry(QtCore.QRect(187, 386, 81, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.leadName_label.setFont(font)
        self.leadName_label.setObjectName("leadName_label")
        self.leadName = QtWidgets.QLineEdit(self.page1)
        self.leadName.setGeometry(QtCore.QRect(260, 390, 301, 31))
        self.leadName.setText("")
        self.leadName.setObjectName("leadName")
        self.leadTitle = QtWidgets.QLineEdit(self.page1)
        self.leadTitle.setGeometry(QtCore.QRect(260, 430, 301, 31))
        self.leadTitle.setObjectName("leadTitle")
        self.eventName = QtWidgets.QLineEdit(self.page1)
        self.eventName.setGeometry(QtCore.QRect(189, 566, 371, 31))
        self.eventName.setText("")
        self.eventName.setObjectName("eventName")
        self.leadTitle_label = QtWidgets.QLabel(self.page1)
        self.leadTitle_label.setGeometry(QtCore.QRect(139, 426, 131, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.leadTitle_label.setFont(font)
        self.leadTitle_label.setObjectName("leadTitle_label")
        self.leadOrg_label = QtWidgets.QLabel(self.page1)
        self.leadOrg_label.setGeometry(QtCore.QRect(110, 467, 151, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.leadOrg_label.setFont(font)
        self.leadOrg_label.setObjectName("leadOrg_label")
        self.endDate = QtWidgets.QDateEdit(self.page1)
        self.endDate.setGeometry(QtCore.QRect(233, 657, 111, 31))
        self.endDate.setDate(QtCore.QDate(2020, 1, 1))
        self.endDate.setObjectName("endDate")
        self.startDate_label = QtWidgets.QLabel(self.page1)
        self.startDate_label.setGeometry(QtCore.QRect(114, 608, 121, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.startDate_label.setFont(font)
        self.startDate_label.setObjectName("startDate_label")
        self.endDate_label = QtWidgets.QLabel(self.page1)
        self.endDate_label.setGeometry(QtCore.QRect(126, 654, 121, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.endDate_label.setFont(font)
        self.endDate_label.setObjectName("endDate_label")
        self.dradisfricData_label = QtWidgets.QLabel(self.page1)
        self.dradisfricData_label.setGeometry(QtCore.QRect(790, 350, 331, 41))
        font = QtGui.QFont()
        font.setPointSize(26)
        font.setBold(True)
        font.setWeight(75)
        self.dradisfricData_label.setFont(font)
        self.dradisfricData_label.setObjectName("dradisfricData_label")
        self.eventName_label = QtWidgets.QLabel(self.page1)
        self.eventName_label.setGeometry(QtCore.QRect(115, 562, 81, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.eventName_label.setFont(font)
        self.eventName_label.setObjectName("eventName_label")
        self.leadOrg = QtWidgets.QLineEdit(self.page1)
        self.leadOrg.setGeometry(QtCore.QRect(260, 471, 301, 31))
        self.leadOrg.setObjectName("leadOrg")
        self.event_label = QtWidgets.QLabel(self.page1)
        self.event_label.setGeometry(QtCore.QRect(80, 508, 131, 41))
        font = QtGui.QFont()
        font.setPointSize(26)
        font.setBold(True)
        font.setWeight(75)
        self.event_label.setFont(font)
        self.event_label.setObjectName("event_label")
        self.teamLead_label = QtWidgets.QLabel(self.page1)
        self.teamLead_label.setGeometry(QtCore.QRect(78, 338, 201, 41))
        font = QtGui.QFont()
        font.setPointSize(26)
        font.setBold(True)
        font.setWeight(75)
        self.teamLead_label.setFont(font)
        self.teamLead_label.setObjectName("teamLead_label")
        self.availFilesFolders_label = QtWidgets.QLabel(self.page1)
        self.availFilesFolders_label.setGeometry(QtCore.QRect(820, 389, 281, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.availFilesFolders_label.setFont(font)
        self.availFilesFolders_label.setObjectName("availFilesFolders_label")
        self.listFolders = QtWidgets.QListWidget(self.page1)
        self.listFolders.setGeometry(QtCore.QRect(810, 420, 371, 321))
        self.listFolders.setObjectName("listFolders")
        self.nextButton = QtWidgets.QPushButton(self.page1)
        self.nextButton.setEnabled(False)
        self.nextButton.setGeometry(QtCore.QRect(1230, 700, 91, 41))
        font = QtGui.QFont()
        font.setPointSize(24)
        font.setBold(True)
        font.setWeight(75)
        self.nextButton.setFont(font)
        self.nextButton.setObjectName("nextButton")
        self.logo_label = QtWidgets.QLabel(self.page1)
        self.logo_label.setGeometry(QtCore.QRect(70, 8, 1161, 321))
        self.logo_label.setText("")
        self.logo_label.setPixmap(QtGui.QPixmap("logos.jpg"))
        self.logo_label.setScaledContents(True)
        self.logo_label.setObjectName("logo_label")
        self.Classificationframe = QtWidgets.QFrame(self.page1)
        self.Classificationframe.setGeometry(QtCore.QRect(600, 537, 191, 181))
        self.Classificationframe.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.Classificationframe.setFrameShadow(QtWidgets.QFrame.Raised)
        self.Classificationframe.setObjectName("Classificationframe")
        self.unclassButton = QtWidgets.QRadioButton(self.Classificationframe)
        self.unclassButton.setGeometry(QtCore.QRect(20, 10, 131, 25))
        self.unclassButton.setChecked(True)
        self.unclassButton.setObjectName("unclassButton")
        self.secretButton = QtWidgets.QRadioButton(self.Classificationframe)
        self.secretButton.setGeometry(QtCore.QRect(20, 86, 81, 25))
        self.secretButton.setObjectName("secretButton")
        self.topsecretButton = QtWidgets.QRadioButton(self.Classificationframe)
        self.topsecretButton.setGeometry(QtCore.QRect(20, 130, 121, 25))
        self.topsecretButton.setObjectName("topsecretButton")
        self.noforncheckBox = QtWidgets.QCheckBox(self.Classificationframe)
        self.noforncheckBox.setEnabled(False)
        self.noforncheckBox.setGeometry(QtCore.QRect(48, 105, 91, 25))
        self.noforncheckBox.setObjectName("noforncheckBox")
        self.scicheckBox = QtWidgets.QCheckBox(self.Classificationframe)
        self.scicheckBox.setEnabled(False)
        self.scicheckBox.setGeometry(QtCore.QRect(48, 149, 51, 25))
        self.scicheckBox.setObjectName("scicheckBox")
        self.frame = QtWidgets.QFrame(self.Classificationframe)
        self.frame.setGeometry(QtCore.QRect(40, 32, 151, 51))
        self.frame.setFrameShape(QtWidgets.QFrame.NoFrame)
        self.frame.setFrameShadow(QtWidgets.QFrame.Plain)
        self.frame.setLineWidth(0)
        self.frame.setObjectName("frame")
        self.fouoButton = QtWidgets.QRadioButton(self.frame)
        self.fouoButton.setGeometry(QtCore.QRect(9, 25, 141, 25))
        self.fouoButton.setObjectName("fouoButton")
        self.cuiButton = QtWidgets.QRadioButton(self.frame)
        self.cuiButton.setGeometry(QtCore.QRect(9, 0, 61, 25))
        self.cuiButton.setChecked(True)
        self.cuiButton.setObjectName("cuiButton")
        self.mode_label = QtWidgets.QLabel(self.page1)
        self.mode_label.setGeometry(QtCore.QRect(450, 607, 111, 51))
        font = QtGui.QFont()
        font.setPointSize(15)
        font.setBold(True)
        font.setWeight(75)
        self.mode_label.setFont(font)
        self.mode_label.setLayoutDirection(QtCore.Qt.LeftToRight)
        self.mode_label.setObjectName("mode_label")
        self.Modeframe = QtWidgets.QFrame(self.page1)
        self.Modeframe.setGeometry(QtCore.QRect(450, 658, 111, 81))
        self.Modeframe.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.Modeframe.setFrameShadow(QtWidgets.QFrame.Raised)
        self.Modeframe.setObjectName("Modeframe")
        self.lightButton = QtWidgets.QRadioButton(self.Modeframe)
        self.lightButton.setGeometry(QtCore.QRect(20, 10, 81, 25))
        self.lightButton.setChecked(False)
        self.lightButton.setObjectName("lightButton")
        self.darkButton = QtWidgets.QRadioButton(self.Modeframe)
        self.darkButton.setGeometry(QtCore.QRect(20, 40, 81, 25))
        self.darkButton.setChecked(True)
        self.darkButton.setObjectName("darkButton")
        self.classification_label = QtWidgets.QLabel(self.page1)
        self.classification_label.setGeometry(QtCore.QRect(604, 506, 151, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.classification_label.setFont(font)
        self.classification_label.setLayoutDirection(QtCore.Qt.LeftToRight)
        self.classification_label.setObjectName("classification_label")
        self.draftcheckBox = QtWidgets.QCheckBox(self.page1)
        self.draftcheckBox.setGeometry(QtCore.QRect(580, 720, 211, 21))
        self.draftcheckBox.setChecked(True)
        self.draftcheckBox.setObjectName("draftcheckBox")
        self.office_comboBox = QtWidgets.QComboBox(self.page1)
        self.office_comboBox.setGeometry(QtCore.QRect(607, 435, 151, 27))
        self.office_comboBox.setEditable(False)
        self.office_comboBox.setCurrentText("")
        self.office_comboBox.setObjectName("office_comboBox")
        sup_list = ['', 'FCDD-DAC-E         (Jai)', 'FCDD-DAC-M         (Isabel)', 'FCDD-DAC-O         (Bert)', 'FCDD-DAC-R         (Justin)']
        self.office_comboBox.addItems(sup_list)
        self.office_label = QtWidgets.QLabel(self.page1)
        self.office_label.setGeometry(QtCore.QRect(609, 400, 151, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.office_label.setFont(font)
        self.office_label.setLayoutDirection(QtCore.Qt.LeftToRight)
        self.office_label.setObjectName("office_label")
        self.Eventframe = QtWidgets.QFrame(self.page1)
        self.Eventframe.setGeometry(QtCore.QRect(196, 698, 201, 41))
        self.Eventframe.setFrameShape(QtWidgets.QFrame.StyledPanel)
        self.Eventframe.setFrameShadow(QtWidgets.QFrame.Raised)
        self.Eventframe.setObjectName("Eventframe")
        self.pmrButton = QtWidgets.QRadioButton(self.Eventframe)
        self.pmrButton.setGeometry(QtCore.QRect(130, 10, 61, 25))
        self.pmrButton.setChecked(False)
        self.pmrButton.setObjectName("pmrButton")
        self.cvpaButton = QtWidgets.QRadioButton(self.Eventframe)
        self.cvpaButton.setGeometry(QtCore.QRect(10, 10, 111, 25))
        self.cvpaButton.setChecked(True)
        self.cvpaButton.setObjectName("cvpaButton")
        self.eventType_label = QtWidgets.QLabel(self.page1)
        self.eventType_label.setGeometry(QtCore.QRect(130, 700, 71, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.eventType_label.setFont(font)
        self.eventType_label.setObjectName("eventType_label")
        self.stackedWidget0.addWidget(self.page1)
        self.page2 = QtWidgets.QWidget()
        self.page2.setObjectName("page2")
        self.quitButton = QtWidgets.QPushButton(self.page2)
        self.quitButton.setGeometry(QtCore.QRect(1220, 640, 91, 41))
        font = QtGui.QFont()
        font.setPointSize(24)
        font.setBold(True)
        font.setWeight(75)
        self.quitButton.setFont(font)
        self.quitButton.setObjectName("quitButton")
        self.findings_label = QtWidgets.QLabel(self.page2)
        self.findings_label.setGeometry(QtCore.QRect(25, 0, 111, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.findings_label.setFont(font)
        self.findings_label.setObjectName("findings_label")
        self.updateDescButton = QtWidgets.QPushButton(self.page2)
        self.updateDescButton.setGeometry(QtCore.QRect(298, 702, 181, 41))
        font = QtGui.QFont()
        font.setPointSize(14)
        font.setBold(True)
        font.setWeight(75)
        self.updateDescButton.setFont(font)
        self.updateDescButton.setObjectName("updateDescButton")
        self.issuesText = QtWidgets.QPlainTextEdit(self.page2)
        self.issuesText.setGeometry(QtCore.QRect(19, 330, 461, 121))
        self.issuesText.setObjectName("issuesText")
        self.screenshotPreview_label = QtWidgets.QLabel(self.page2)
        self.screenshotPreview_label.setGeometry(QtCore.QRect(550, 352, 231, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.screenshotPreview_label.setFont(font)
        self.screenshotPreview_label.setObjectName("screenshotPreview_label")
        self.imagePreview_label = QtWidgets.QLabel(self.page2)
        self.imagePreview_label.setGeometry(QtCore.QRect(540, 383, 631, 391))
        self.imagePreview_label.setScaledContents(True)
        self.imagePreview_label.setAlignment(QtCore.Qt.AlignCenter)
        self.imagePreview_label.setObjectName("imagePreview_label")
        self.findingDeleteButton = QtWidgets.QPushButton(self.page2)
        self.findingDeleteButton.setGeometry(QtCore.QRect(439, 178, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.findingDeleteButton.setFont(font)
        self.findingDeleteButton.setObjectName("findingDeleteButton")
        self.findingUpButton = QtWidgets.QPushButton(self.page2)
        self.findingUpButton.setGeometry(QtCore.QRect(440, 98, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.findingUpButton.setFont(font)
        self.findingUpButton.setCursor(QtGui.QCursor(QtCore.Qt.ArrowCursor))
        self.findingUpButton.setObjectName("findingUpButton")
        self.findingDownButton = QtWidgets.QPushButton(self.page2)
        self.findingDownButton.setGeometry(QtCore.QRect(440, 128, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.findingDownButton.setFont(font)
        self.findingDownButton.setObjectName("findingDownButton")
        self.screenshots_label = QtWidgets.QLabel(self.page2)
        self.screenshots_label.setGeometry(QtCore.QRect(566, 0, 141, 31))
        font = QtGui.QFont()
        font.setPointSize(18)
        font.setBold(True)
        font.setWeight(75)
        self.screenshots_label.setFont(font)
        self.screenshots_label.setObjectName("screenshots_label")
        self.pptxButton = QtWidgets.QPushButton(self.page2)
        self.pptxButton.setGeometry(QtCore.QRect(1220, 30, 91, 41))
        font = QtGui.QFont()
        font.setPointSize(24)
        font.setBold(True)
        font.setWeight(75)
        self.pptxButton.setFont(font)
        self.pptxButton.setObjectName("pptxButton")
        self.issues_label = QtWidgets.QLabel(self.page2)
        self.issues_label.setGeometry(QtCore.QRect(29, 299, 51, 31))
        font = QtGui.QFont()
        font.setPointSize(11)
        font.setBold(True)
        font.setWeight(75)
        self.issues_label.setFont(font)
        self.issues_label.setObjectName("issues_label")
        self.listFindings = QtWidgets.QListWidget(self.page2)
        self.listFindings.setGeometry(QtCore.QRect(15, 30, 421, 181))
        self.listFindings.setDragEnabled(False)
        self.listFindings.setDragDropOverwriteMode(False)
        self.listFindings.setDragDropMode(QtWidgets.QAbstractItemView.NoDragDrop)
        self.listFindings.setDefaultDropAction(QtCore.Qt.IgnoreAction)
        self.listFindings.setObjectName("listFindings")
        self.gobackButton = QtWidgets.QPushButton(self.page2)
        self.gobackButton.setGeometry(QtCore.QRect(1200, 700, 131, 41))
        font = QtGui.QFont()
        font.setPointSize(24)
        font.setBold(True)
        font.setWeight(75)
        self.gobackButton.setFont(font)
        self.gobackButton.setObjectName("gobackButton")
        self.screenshotDeleteButton = QtWidgets.QPushButton(self.page2)
        self.screenshotDeleteButton.setGeometry(QtCore.QRect(1129, 240, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.screenshotDeleteButton.setFont(font)
        self.screenshotDeleteButton.setObjectName("screenshotDeleteButton")
        self.screenshotDownButton = QtWidgets.QPushButton(self.page2)
        self.screenshotDownButton.setGeometry(QtCore.QRect(1130, 190, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.screenshotDownButton.setFont(font)
        self.screenshotDownButton.setObjectName("screenshotDownButton")
        self.screenshotUpButton = QtWidgets.QPushButton(self.page2)
        self.screenshotUpButton.setGeometry(QtCore.QRect(1130, 160, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.screenshotUpButton.setFont(font)
        self.screenshotUpButton.setCursor(QtGui.QCursor(QtCore.Qt.ArrowCursor))
        self.screenshotUpButton.setObjectName("screenshotUpButton")
        self.listScreenshots = QtWidgets.QListWidget(self.page2)
        self.listScreenshots.setGeometry(QtCore.QRect(556, 30, 571, 321))
        self.listScreenshots.setDragEnabled(False)
        self.listScreenshots.setDragDropOverwriteMode(False)
        self.listScreenshots.setDragDropMode(QtWidgets.QAbstractItemView.NoDragDrop)
        self.listScreenshots.setDefaultDropAction(QtCore.Qt.IgnoreAction)
        self.listScreenshots.setObjectName("listScreenshots")
        self.findingName_label = QtWidgets.QLabel(self.page2)
        self.findingName_label.setGeometry(QtCore.QRect(33, 228, 101, 31))
        font = QtGui.QFont()
        font.setPointSize(11)
        font.setBold(True)
        font.setWeight(75)
        self.findingName_label.setFont(font)
        self.findingName_label.setObjectName("findingName_label")
        self.findingName = QtWidgets.QLineEdit(self.page2)
        self.findingName.setGeometry(QtCore.QRect(139, 229, 341, 31))
        self.findingName.setText("")
        self.findingName.setPlaceholderText("")
        self.findingName.setObjectName("findingName")
        self.findingHosts_label = QtWidgets.QLabel(self.page2)
        self.findingHosts_label.setGeometry(QtCore.QRect(29, 267, 111, 31))
        font = QtGui.QFont()
        font.setPointSize(11)
        font.setBold(True)
        font.setWeight(75)
        self.findingHosts_label.setFont(font)
        self.findingHosts_label.setObjectName("findingHosts_label")
        self.findingHosts = QtWidgets.QLineEdit(self.page2)
        self.findingHosts.setGeometry(QtCore.QRect(139, 269, 341, 31))
        self.findingHosts.setText("")
        self.findingHosts.setPlaceholderText("")
        self.findingHosts.setObjectName("findingHosts")
        self.posture_label = QtWidgets.QLabel(self.page2)
        self.posture_label.setGeometry(QtCore.QRect(19, 450, 61, 31))
        font = QtGui.QFont()
        font.setPointSize(11)
        font.setBold(True)
        font.setWeight(75)
        self.posture_label.setFont(font)
        self.posture_label.setObjectName("posture_label")
        self.insiderButton = QtWidgets.QRadioButton(self.page2)
        self.insiderButton.setGeometry(QtCore.QRect(36, 477, 91, 25))
        self.insiderButton.setObjectName("insiderButton")
        self.nearsiderButton = QtWidgets.QRadioButton(self.page2)
        self.nearsiderButton.setGeometry(QtCore.QRect(36, 497, 91, 25))
        self.nearsiderButton.setChecked(True)
        self.nearsiderButton.setObjectName("nearsiderButton")
        self.outsiderButton = QtWidgets.QRadioButton(self.page2)
        self.outsiderButton.setGeometry(QtCore.QRect(36, 517, 91, 25))
        self.outsiderButton.setChecked(False)
        self.outsiderButton.setObjectName("outsiderButton")
        self.mitigationText = QtWidgets.QPlainTextEdit(self.page2)
        self.mitigationText.setGeometry(QtCore.QRect(19, 566, 461, 121))
        self.mitigationText.setObjectName("mitigationText")
        self.mitigation_label = QtWidgets.QLabel(self.page2)
        self.mitigation_label.setGeometry(QtCore.QRect(29, 539, 81, 31))
        font = QtGui.QFont()
        font.setPointSize(11)
        font.setBold(True)
        font.setWeight(75)
        self.mitigation_label.setFont(font)
        self.mitigation_label.setObjectName("mitigation_label")
        self.mitigationcheckBox = QtWidgets.QCheckBox(self.page2)
        self.mitigationcheckBox.setGeometry(QtCore.QRect(29, 696, 141, 25))
        self.mitigationcheckBox.setObjectName("mitigationcheckBox")
        self.findingAddButton = QtWidgets.QPushButton(self.page2)
        self.findingAddButton.setGeometry(QtCore.QRect(440, 48, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.findingAddButton.setFont(font)
        self.findingAddButton.setObjectName("findingAddButton")
        self.screenshotAddButton = QtWidgets.QPushButton(self.page2)
        self.screenshotAddButton.setGeometry(QtCore.QRect(1130, 100, 41, 21))
        font = QtGui.QFont()
        font.setPointSize(8)
        font.setBold(True)
        font.setWeight(75)
        self.screenshotAddButton.setFont(font)
        self.screenshotAddButton.setObjectName("screenshotAddButton")
        self.imagePreview_label.raise_()
        self.quitButton.raise_()
        self.findings_label.raise_()
        self.updateDescButton.raise_()
        self.issuesText.raise_()
        self.screenshotPreview_label.raise_()
        self.findingDeleteButton.raise_()
        self.findingUpButton.raise_()
        self.findingDownButton.raise_()
        self.screenshots_label.raise_()
        self.pptxButton.raise_()
        self.issues_label.raise_()
        self.listFindings.raise_()
        self.gobackButton.raise_()
        self.screenshotDeleteButton.raise_()
        self.screenshotDownButton.raise_()
        self.screenshotUpButton.raise_()
        self.listScreenshots.raise_()
        self.findingName_label.raise_()
        self.findingName.raise_()
        self.findingHosts_label.raise_()
        self.findingHosts.raise_()
        self.posture_label.raise_()
        self.insiderButton.raise_()
        self.nearsiderButton.raise_()
        self.outsiderButton.raise_()
        self.mitigationText.raise_()
        self.mitigation_label.raise_()
        self.mitigationcheckBox.raise_()
        self.findingAddButton.raise_()
        self.screenshotAddButton.raise_()
        self.stackedWidget0.addWidget(self.page2)
        MainWindow.setCentralWidget(self.centralwidget)
        self.menubar = QtWidgets.QMenuBar(MainWindow)
        self.menubar.setGeometry(QtCore.QRect(0, 0, 1345, 24))
        self.menubar.setObjectName("menubar")
        MainWindow.setMenuBar(self.menubar)
        self.statusbar = QtWidgets.QStatusBar(MainWindow)
        self.statusbar.setObjectName("statusbar")
        MainWindow.setStatusBar(self.statusbar)

        self.retranslateUi(MainWindow)
        self.stackedWidget0.setCurrentIndex(0)
        QtCore.QMetaObject.connectSlotsByName(MainWindow)
        MainWindow.setTabOrder(self.leadName, self.leadTitle)
        MainWindow.setTabOrder(self.leadTitle, self.leadOrg)
        MainWindow.setTabOrder(self.leadOrg, self.eventName)
        MainWindow.setTabOrder(self.eventName, self.startDate)
        MainWindow.setTabOrder(self.startDate, self.endDate)
        MainWindow.setTabOrder(self.endDate, self.cvpaButton)
        MainWindow.setTabOrder(self.cvpaButton, self.pmrButton)
        MainWindow.setTabOrder(self.pmrButton, self.lightButton)
        MainWindow.setTabOrder(self.lightButton, self.darkButton)
        MainWindow.setTabOrder(self.darkButton, self.office_comboBox)
        MainWindow.setTabOrder(self.office_comboBox, self.unclassButton)
        MainWindow.setTabOrder(self.unclassButton, self.cuiButton)
        MainWindow.setTabOrder(self.cuiButton, self.fouoButton)
        MainWindow.setTabOrder(self.fouoButton, self.secretButton)
        MainWindow.setTabOrder(self.secretButton, self.noforncheckBox)
        MainWindow.setTabOrder(self.noforncheckBox, self.topsecretButton)
        MainWindow.setTabOrder(self.topsecretButton, self.scicheckBox)
        MainWindow.setTabOrder(self.scicheckBox, self.draftcheckBox)
        MainWindow.setTabOrder(self.draftcheckBox, self.listFolders)
        MainWindow.setTabOrder(self.listFolders, self.nextButton)
        MainWindow.setTabOrder(self.nextButton, self.listFindings)
        MainWindow.setTabOrder(self.listFindings, self.findingAddButton)
        MainWindow.setTabOrder(self.findingAddButton, self.findingUpButton)
        MainWindow.setTabOrder(self.findingUpButton, self.findingDownButton)
        MainWindow.setTabOrder(self.findingDownButton, self.findingDeleteButton)
        MainWindow.setTabOrder(self.findingDeleteButton, self.findingName)
        MainWindow.setTabOrder(self.findingName, self.findingHosts)
        MainWindow.setTabOrder(self.findingHosts, self.issuesText)
        MainWindow.setTabOrder(self.issuesText, self.insiderButton)
        MainWindow.setTabOrder(self.insiderButton, self.nearsiderButton)
        MainWindow.setTabOrder(self.nearsiderButton, self.outsiderButton)
        MainWindow.setTabOrder(self.outsiderButton, self.mitigationText)
        MainWindow.setTabOrder(self.mitigationText, self.mitigationcheckBox)
        MainWindow.setTabOrder(self.mitigationcheckBox, self.updateDescButton)
        MainWindow.setTabOrder(self.updateDescButton, self.listScreenshots)
        MainWindow.setTabOrder(self.listScreenshots, self.screenshotAddButton)
        MainWindow.setTabOrder(self.screenshotAddButton, self.screenshotUpButton)
        MainWindow.setTabOrder(self.screenshotUpButton, self.screenshotDownButton)
        MainWindow.setTabOrder(self.screenshotDownButton, self.screenshotDeleteButton)
        MainWindow.setTabOrder(self.screenshotDeleteButton, self.pptxButton)
        MainWindow.setTabOrder(self.pptxButton, self.gobackButton)
        MainWindow.setTabOrder(self.gobackButton, self.quitButton)


        class Event:
            def __init__(self, lead_name, lead_title, lead_org, office_symbol, event_name, event_type, start_date, end_date, classification, designation, draft, mode):
                self.lead_name = lead_name
                self.lead_title = lead_title
                self.lead_org = lead_org
                self.office_symbol = office_symbol
                self.event_name = event_name
                self.event_type = event_type
                self.start_date = start_date
                self.end_date = end_date
                self.classification = classification
                self.designation = designation
                self.draft = draft
                self.mode = mode

            def write_file(self):
                #CREATE NEW XML FILE
                xml_file = cwd + '/event.xml'
                def indent(elem, level=0):
                    i = "\n" + level*"    "
                    if len(elem):
                        if not elem.text or not elem.text.strip():
                            elem.text = i + "    "
                        if not elem.tail or not elem.tail.strip():
                            elem.tail = i
                        for elem in elem:
                            indent(elem, level+1)
                        if not elem.tail or not elem.tail.strip():
                            elem.tail = i
                    else:
                        if level and (not elem.tail or not elem.tail.strip()):
                            elem.tail = i
                #CREATE FILE STRUCTURE
                root = ET.Element('data')
                #CREATE EVENT
                xml_event = ET.SubElement(root, 'event')
                xml_event.set('uid', str(0))
                xml_lead_name = ET.SubElement(xml_event, 'lead_name')
                xml_lead_name.text = self.lead_name
                xml_lead_title = ET.SubElement(xml_event, 'lead_title')
                xml_lead_title.text = self.lead_title
                xml_lead_org = ET.SubElement(xml_event, 'lead_org')
                xml_lead_org.text = self.lead_org
                xml_office_symbol = ET.SubElement(xml_event, 'office_symbol')
                xml_office_symbol.text = self.office_symbol
                xml_event_name = ET.SubElement(xml_event, 'event_name')
                xml_event_name.text = self.event_name
                xml_event_type = ET.SubElement(xml_event, 'event_type')
                xml_event_type.text = self.event_type
                xml_start_date = ET.SubElement(xml_event, 'start_date')
                xml_start_date.text = self.start_date
                xml_end_date = ET.SubElement(xml_event, 'end_date')
                xml_end_date.text = self.end_date
                xml_classification = ET.SubElement(xml_event, 'classification')
                xml_classification.text = self.classification
                xml_designation = ET.SubElement(xml_event, 'designation')
                xml_designation.text = self.designation
                xml_draft = ET.SubElement(xml_event, 'draft')
                xml_draft.text = self.draft
                xml_mode = ET.SubElement(xml_event, 'mode')
                xml_mode.text = self.mode
                #WRITING XML
                indent(root)
                tree = ET.ElementTree(root)
                tree.write(xml_file, encoding='utf-8', xml_declaration=True)


        class Finding:
            def __init__(self, folder, active, rank, title, hosts, issues, posture, mitigation, include_mitigation, screenshots=None):
                self.folder = folder
                self.active = active
                self.rank = rank
                self.title = title
                self.hosts = hosts
                self.issues = issues
                self.posture = posture
                self.mitigation = mitigation
                self.include_mitigation = include_mitigation
                self.screenshots = screenshots

            def insert_new_finding(self, uid):
                def indent(elem, level=0):
                    i = "\n" + level*"    "
                    if len(elem):
                        if not elem.text or not elem.text.strip():
                            elem.text = i + "    "
                        if not elem.tail or not elem.tail.strip():
                            elem.tail = i
                        for elem in elem:
                            indent(elem, level+1)
                        if not elem.tail or not elem.tail.strip():
                            elem.tail = i
                    else:
                        if level and (not elem.tail or not elem.tail.strip()):
                            elem.tail = i
                #UPDATE XML FILE
                tree = ET.parse(current_erb)
                root = tree.getroot()
                xml_finding = ET.SubElement(root, 'finding')
                xml_finding.set('uid', uid)
                xml_folder = ET.SubElement(xml_finding, 'folder')
                xml_folder.text = str(self.folder)
                xml_active = ET.SubElement(xml_finding, 'active')
                xml_active.text = str(self.active)
                xml_rank = ET.SubElement(xml_finding, 'rank')
                xml_rank.text = str(self.rank)
                xml_title = ET.SubElement(xml_finding, 'title')
                xml_title.text = str(self.title)
                xml_hosts = ET.SubElement(xml_finding, 'hosts')
                xml_hosts.text = str(self.hosts)
                xml_issues = ET.SubElement(xml_finding, 'issues')
                xml_issues.text = str(self.issues)
                xml_posture = ET.SubElement(xml_finding, 'posture')
                xml_posture.text = str(self.posture)
                xml_mitigation = ET.SubElement(xml_finding, 'mitigation')
                xml_mitigation.text = str(self.mitigation)
                xml_include_mitigation = ET.SubElement(xml_finding, 'include_mitigation')
                xml_include_mitigation.text = str(self.include_mitigation)
                xml_screenshots = ET.SubElement(xml_finding, 'screenshots')
                xml_screenshots.text = str(self.screenshots)
                #WRITING XML
                indent(root)
                tree.write(current_erb, encoding='utf-8', xml_declaration=True)

            def update_attributes(self, finding_num, title, hosts, issues, posture, mitigation, include_mitigation):
                self.title = title
                self.hosts = hosts
                self.issues = issues
                self.posture = posture
                self.mitigation = mitigation
                self.include_mitigation = include_mitigation
                #UPDATE FINDING IN XML
                tree = ET.parse(current_erb)
                root = tree.getroot()
                for finding in root.findall('finding'):
                    uid = finding.get('uid')
                    if uid == finding_num:
                        xml_title = finding.find('title')
                        xml_hosts = finding.find('hosts')
                        xml_issues = finding.find('issues')
                        xml_posture = finding.find('posture')
                        xml_mitigation = finding.find('mitigation')
                        xml_include_mitigation = finding.find('include_mitigation')
                        xml_title.text = self.title
                        xml_hosts.text = self.hosts
                        xml_issues.text = self.issues
                        xml_posture.text = self.posture
                        xml_mitigation.text = self.mitigation
                        xml_include_mitigation.text = self.include_mitigation
                tree.write(current_erb)

            @classmethod
            def modify_rank(self):
                #MODIFY RANK ORDER IN XML
                tree = ET.parse(current_erb)
                root = tree.getroot()
                for finding in root.findall('finding'):
                    uid = finding.get('uid')
                    active = finding.find('active').text
                    if active == '1':
                        rank = finding.find('rank')
                        value = order_list.index(uid)
                        rank.text = str(value)
                tree.write(current_erb)

            @classmethod
            def deactivate(self, finding_num):
                #DE-ACTIVATE FINDING IN XML
                tree = ET.parse(current_erb)
                root = tree.getroot()
                for finding in root.findall('finding'):
                    uid = finding.get('uid')
                    if uid == order_list[finding_num]:
                        active = finding.find('active')
                        rank = finding.find('rank')
                        active.text = '0'
                        rank.text = 'x'
                tree.write(current_erb)
                order_list.remove(order_list[finding_num])

            def modify_folder(self, finding_num):
                #MODIFY SCREENSHOTS FOLDER IN XML
                tree = ET.parse(current_erb)
                root = tree.getroot()
                for finding in root.findall('finding'):
                    uid = finding.get('uid')
                    if uid == order_list[finding_num]:
                        screenshots = finding.find('folder')
                        screenshots.text = str(self.folder)
                tree.write(current_erb)

            def modify_screenshots(self, finding_num):
                #MODIFY SCREENSHOTS ORDER IN XML
                tree = ET.parse(current_erb)
                root = tree.getroot()
                for finding in root.findall('finding'):
                    uid = finding.get('uid')
                    if uid == order_list[finding_num]:
                        screenshots = finding.find('screenshots')
                        screenshots.text = str(self.screenshots)
                tree.write(current_erb)


        #CLEAR ALL FIELDS ON PAGE 2
        def clear_all():
            self.listFindings.clear()
            self.findingName.setText("")
            self.findingHosts.setText("")
            self.issuesText.clear()
            if self.nearsiderButton.isChecked() == False:
                self.nearsiderButton.toggle()
            self.mitigationText.clear()
            if self.mitigationcheckBox.isChecked() == True:
                self.mitigationcheckBox.toggle()
            self.listScreenshots.clear()
            self.imagePreview_label.setText("NO PREVIEW")


        #DISPLAY IMAGES ON "SCREENSHOT PREVIEW"
        def preview_image():
            if self.listScreenshots.selectedItems() != []:
                item1 = self.listFindings.selectedIndexes()[0]
                finding_num = item1.row()
                current_finding = order_list[finding_num]
                item2 = self.listScreenshots.selectedIndexes()[0]
                image_num = item2.row()
                image_name = findings[current_finding].screenshots[image_num]
                image_file = findings[str(order_list[finding_num])].folder + image_name
                self.imagePreview_label.setText("")
                self.imagePreview_label.setPixmap(QtGui.QPixmap(image_file))
                self.imagePreview_label.setScaledContents(True)
            else:
                self.imagePreview_label.setText("NO PREVIEW")


        #POPULATE SCREENSHOTS LIST
        def set_screenshots():
            self.listScreenshots.clear()
            item = self.listFindings.selectedIndexes()[0]
            index = item.row()
            current_finding = order_list[index]
            if findings[current_finding].screenshots != None:
                for i in findings[current_finding].screenshots:
                    self.listScreenshots.addItem(i)
                if findings[current_finding].screenshots != [] and self.listScreenshots.selectedItems() == []:
                    self.listScreenshots.item(0).setSelected(True)
                    self.listScreenshots.setCurrentRow(0)


        #POPULATE FINDINGS FIELDS
        def set_fields():
            if self.listFindings.selectedItems() != []:
                item = self.listFindings.selectedIndexes()[0]
                index = item.row()
                current_finding = order_list[index]
                self.findingName.setText(findings[current_finding].title)
                self.findingHosts.setText(findings[current_finding].hosts)
                self.issuesText.clear()
                self.issuesText.insertPlainText(findings[current_finding].issues)
                if findings[current_finding].posture == 'INSIDER':
                    if self.insiderButton.isChecked() == False:
                        self.insiderButton.toggle()
                if findings[current_finding].posture == 'NEARSIDER':
                    if self.nearsiderButton.isChecked() == False:
                        self.nearsiderButton.toggle()
                if findings[current_finding].posture == 'OUTSIDER':
                    if self.outsiderButton.isChecked() == False:
                        self.outsiderButton.toggle()
                self.mitigationText.clear()
                self.mitigationText.insertPlainText(findings[current_finding].mitigation)
                if findings[current_finding].include_mitigation == 'yes':
                    if self.mitigationcheckBox.isChecked() == False:
                        self.mitigationcheckBox.toggle()
                if findings[current_finding].include_mitigation == 'no':
                    if self.mitigationcheckBox.isChecked() == True:
                        self.mitigationcheckBox.toggle()
                set_screenshots()


        #UPDATE FINDING XML
        def update_finding():
            if self.listFindings.selectedItems() != []:
                item = self.listFindings.selectedIndexes()[0]
                index = item.row()
                current_finding = order_list[index]
                title = self.findingName.text()
                hosts = self.findingHosts.text()
                issues = self.issuesText.toPlainText()
                if self.insiderButton.isChecked() == True:
                    posture = 'INSIDER'
                if self.nearsiderButton.isChecked() == True:
                    posture = 'NEARSIDER'
                if self.outsiderButton.isChecked() == True:
                    posture = 'OUTSIDER'
                mitigation = self.mitigationText.toPlainText()
                if self.mitigationcheckBox.isChecked() == True:
                    include_mitigation = 'yes'
                else:
                    include_mitigation = 'no'
                #UPDATE FINDING
                findings[str(current_finding)].update_attributes(current_finding, title, hosts, issues, posture, mitigation, include_mitigation)
                self.listFindings.clear()
                for i in order_list:
                    self.listFindings.addItem(findings[i].title)
                self.listFindings.item(index).setSelected(True)
                self.listFindings.setCurrentRow(index)
                #MESSAGE BOX!
                msg = QMessageBox()
                msg.setIcon(QMessageBox.Information)
                msg.setWindowTitle("Update Finding")
                msg.setText("Finding was successfully updated!")
                x = msg.exec_()


        #GET AVAILABLE DRADIS/FRIC FILES/FOLDERS FROM "DESKTOP" + CURRENT ERB
        def get_folders():
            global cwd, current_erb
            if os.path.isfile(current_erb):
                self.listFolders.addItem('Existing ERB')
            self.listFolders.addItem('Create your own ERB')
            os.chdir(desktop_dir)
            for file in glob.glob("dradis-export*.zip"):
                self.listFolders.addItem(str(file))
            raw_folders = next(os.walk(desktop_dir))[1]
            prefix = "fric_export_"
            for i in raw_folders:
                if prefix in i:
                    self.listFolders.addItem(str(i))
            os.chdir(cwd)
        get_folders()


        #GET FINDINGS FROM DRADIS/FRIC FILES/FOLDERS
        def get_findings():
            global order_list
            tree = ET.parse(current_erb)
            root = tree.getroot()
            num_findings = len(root)
            findings.clear()
            self.listFindings.clear()
            uids = []
            ranks = []
            del order_list[:]
            for finding in root.findall('finding'):
                uid = finding.get('uid')
                folder = finding.find('folder').text
                active = finding.find('active').text
                rank = finding.find('rank').text
                title = finding.find('title').text
                hosts = finding.find('hosts').text
                issues = finding.find('issues').text
                posture = finding.find('posture').text
                mitigation = finding.find('mitigation').text
                include_mitigation = finding.find('include_mitigation').text
                tmp = finding.find('screenshots').text
                if tmp != None:
                    tmp = tmp.split('[')
                    tmp = tmp[1].split(']')
                    tmp = tmp[0]
                    tmp = tmp.split(', ')
                    for k in range(len(tmp)):
                        tmp2 = tmp[k].split("'")
                        if len(tmp2) > 1:
                            tmp[k] = tmp2[1]
                            screenshots = tmp
                        else:
                            screenshots = None
                else:
                    screenshots = None
                if active == '0':
                    rank = 'x'
                else:
                    findings[uid] = Finding(folder, active, rank, title, hosts, issues, posture, mitigation, include_mitigation, screenshots)
                    uids.append(uid)
                    ranks.append(rank)
            int_ranks = [int(x) for x in ranks]
            order_list = [x for _,x in sorted(zip(int_ranks,uids))]
            #POPULATE FINDINGS LIST
            for i in order_list:
                self.listFindings.addItem(findings[i].title)


        #CHECK IF ITEM IS SELECTED FROM DRADIS/FRIC FILES/FOLDERS
        def on_selection_changed():
            clear_all()
            global data_folder, data_source
            data_source = self.listFolders.currentItem().text()
            if self.listFolders.selectedItems():
                self.nextButton.setEnabled(True)
                if 'Existing ERB' in data_source:
                    data_folder = cwd + '/erb/'
                elif 'Create your own ERB' in data_source:
                    data_folder = cwd + '/erb/'
                else:
                    data_folder = desktop_dir + '/' + data_source


        #CHECK IF EVENT XML FILE EXISTS
        xml_file = cwd + '/event.xml'
        if os.path.isfile(xml_file):
            tree = ET.parse(xml_file)
            root = tree.getroot()
            for event in root.findall('event'):
                uid = event.get('uid')
                lead_name = event.find('lead_name').text
                lead_title = event.find('lead_title').text
                lead_org = event.find('lead_org').text
                office_symbol = event.find('office_symbol').text
                event_name = event.find('event_name').text
                event_type = event.find('event_type').text
                start_date = event.find('start_date').text
                end_date = event.find('end_date').text
                classification = event.find('classification').text
                designation = event.find('designation').text
                draft = event.find('draft').text
                mode = event.find('mode').text
                
                self.leadName.setText(lead_name)
                self.leadTitle.setText(lead_title)
                self.leadOrg.setText(lead_org)
                self.office_comboBox.setCurrentText(office_symbol)
                self.eventName.setText(event_name)
                if event_type == 'CVPA':
                    self.cvpaButton.setChecked(True)
                else:
                    self.pmrButton.setChecked(True)
                    self.dradisfricData_label.hide()
                    self.availFilesFolders_label.hide()
                    self.listFolders.hide()
                    self.nextButton.setText("PPTX")
                    self.nextButton.setEnabled(True)
                start_date = start_date.split('/')
                self.startDate.setDate(QDate(int(start_date[2]), int(start_date[0]), int(start_date[1])))
                end_date = end_date.split('/')
                self.endDate.setDate(QDate(int(end_date[2]), int(end_date[0]), int(end_date[1])))
                self.cuiButton.setEnabled(False)
                self.fouoButton.setEnabled(False)
                if classification == 'UNCLASSIFIED':
                    self.unclassButton.setChecked(True)
                    self.cuiButton.setEnabled(True)
                    self.fouoButton.setEnabled(True)
                    if designation == 'CUI':
                        self.cuiButton.setChecked(True)
                    elif designation == 'FOUO':
                        self.fouoButton.setChecked(True)
                    else:
                        self.fouoButton.setChecked(False)
                if classification == 'SECRET':
                    self.secretButton.setChecked(True)
                    self.noforncheckBox.setEnabled(True)
                    if designation == 'NOFORN':
                        self.noforncheckBox.setChecked(True)
                if classification == 'TOP SECRET':
                    self.topsecretButton.setChecked(True)
                    self.scicheckBox.setEnabled(True)
                    if designation == 'SCI':
                        self.scicheckBox.setChecked(True)
                if draft == 'YES':
                    self.draftcheckBox.setChecked(True)
                else:
                    self.draftcheckBox.setChecked(False)
                if mode == 'LIGHT':
                    if self.lightButton.isChecked() == False:
                        self.lightButton.toggle()
                if mode == 'DARK':
                    if self.darkButton.isChecked() == False:
                        self.darkButton.toggle()
        else:
            #HANDLE DATES
            #TODAY'S DATE
            now = QDate.currentDate()
            #TODAY'S WEEK DAY - GETTING INDEX/OFFSET
            week_days = ['Mon', 'Tue', 'Wed', 'Thu', 'Fri', 'Sat', 'Sun']
            datetime = QDateTime.currentDateTime()
            datetime = datetime.toString()
            week_day = datetime[0:3]
            todays_index = week_days.index(week_day) #Mon - 0, Tue - 1, ..., Sun - 6
            #DETERMINE START DAY
            start_day = now.addDays(-todays_index)
            #DETERMINE END DAY
            end_day = start_day.addDays(4)
            #SETTING NEW DATES - START / END
            self.startDate.setDate(start_day)
            self.endDate.setDate(end_day)


        #EVENT TYPE SELECTION
        def event_type():
            if self.cvpaButton.isChecked() == True:
                self.dradisfricData_label.show()
                self.availFilesFolders_label.show()
                self.listFolders.show()
                self.nextButton.setText("Next")
                if self.listFolders.selectedItems():
                    self.nextButton.setEnabled(True)
                else:
                    self.nextButton.setEnabled(False)
            if self.pmrButton.isChecked() == True:
                self.dradisfricData_label.hide()
                self.availFilesFolders_label.hide()
                self.listFolders.hide()
                self.nextButton.setText("PPTX")
                self.nextButton.setEnabled(True)


        #CLASSIFICATION SELECTION CHANGE
        def class_changed():
            self.noforncheckBox.setChecked(False)
            self.scicheckBox.setChecked(False)
            if self.unclassButton.isChecked():
                self.cuiButton.setEnabled(True)
                self.fouoButton.setEnabled(True)
                self.noforncheckBox.setEnabled(False)
                self.scicheckBox.setEnabled(False)
            if self.secretButton.isChecked():
                self.cuiButton.setEnabled(False)
                self.fouoButton.setEnabled(False)
                self.cuiButton.setChecked(False)
                self.fouoButton.setChecked(False)
                self.noforncheckBox.setEnabled(True)
                self.scicheckBox.setEnabled(False)
            if self.topsecretButton.isChecked():
                self.cuiButton.setEnabled(False)
                self.fouoButton.setEnabled(False)
                self.cuiButton.setChecked(False)
                self.fouoButton.setChecked(False)
                self.noforncheckBox.setEnabled(False)
                self.scicheckBox.setEnabled(True)


	#ADD FINDING TO LIST
        def addFinding():
            global order_list
            uid = str(len(findings))
            folder = ''
            active = 1
            rank = str(len(findings))
            title = 'NEW FINDING'
            hosts = 'ADD HOST(S)'
            issues = 'ADD ISSUE(S)'
            posture = 'NEARSIDER'
            mitigation = 'NO MITIGATION FOUND FOR THIS FINDING...'
            include_mitigation = 'yes'
            screenshots = []
            findings[uid] = Finding(folder, active, rank, title, hosts, issues, posture, mitigation, include_mitigation, screenshots)
            findings[uid].insert_new_finding(uid)
            order_list.append(uid)
            self.listFindings.addItem(title)
            current_list = [str(self.listFindings.item(i).text()) for i in range(self.listFindings.count())]
            index = len(current_list)-1
            self.listFindings.item(index).setSelected(True)
            self.listFindings.setCurrentRow(index)


        #MOVE ONE FINDING UP THE LIST
        def findingsUp():
            global order_list
            currentRow = self.listFindings.currentRow()
            currentItem = self.listFindings.takeItem(currentRow)
            if currentRow > 0:
                tmp = order_list[currentRow-1]
                order_list[currentRow-1] = order_list[currentRow]
                order_list[currentRow] = tmp
            if currentRow != -1:
                self.listFindings.insertItem(currentRow - 1, currentItem)
                self.listFindings.setCurrentItem(currentItem)
                selected_finding = self.listFindings.currentItem().text()
                current_list = [str(self.listFindings.item(i).text()) for i in range(self.listFindings.count())]
                Finding.modify_rank()


        #MOVE ONE FINDING DOWN THE LIST
        def findingsDown():
            global order_list
            currentRow = self.listFindings.currentRow()
            currentItem = self.listFindings.takeItem(currentRow)
            if currentRow < len(order_list)-1:
                tmp = order_list[currentRow+1]
                order_list[currentRow+1] = order_list[currentRow]
                order_list[currentRow] = tmp
            if currentRow != -1:
                self.listFindings.insertItem(currentRow + 1, currentItem)
                self.listFindings.setCurrentItem(currentItem)
                selected_finding = self.listFindings.currentItem().text()
                current_list = [str(self.listFindings.item(i).text()) for i in range(self.listFindings.count())]
                Finding.modify_rank()


        #DELETE FINDING FROM LIST
        def delFinding():
            global order_list
            currentRow = self.listFindings.currentRow()
            item = self.listFindings.takeItem(self.listFindings.currentRow())
            item = None
            current_list = [str(self.listFindings.item(i).text()) for i in range(self.listFindings.count())]
            if currentRow != -1:
                Finding.deactivate(currentRow)
            if current_list == []:
                clear_all()


        #ADD SCREENSHOT TO LIST
        def addScreenshot():
            global f_folder, screenshots
            if self.listFindings.selectedItems() != []:
                currentRow = self.listFindings.currentRow()
                if f_folder == '':
                    f_folder = desktop_dir
                fname, _filter = QtWidgets.QFileDialog.getOpenFileName(None, "Import Image", f_folder, "*.png *.jpg *.bmp *.gif *.jpeg *.tif *.tiff")
                if fname != '':
                    aux = fname.rsplit('/', 1)
                    f_folder = aux[0]
                    new_image = fname.split('/')
                    new_image = new_image[-1]
                    findings[order_list[currentRow]].screenshots.append(new_image)
                    if findings[order_list[currentRow]].folder == '':
                        sf = cwd + '/erb/'
                        num_folders = len(next(os.walk(sf))[1])
                        img_folder = sf + str(num_folders)
                        msg = 'mkdir ' + img_folder
                        os.system(msg)
                        findings[order_list[currentRow]].folder = img_folder + '/'
                        findings[order_list[currentRow]].modify_folder(currentRow)
                    else:
                        img_folder = findings[order_list[currentRow]].folder
                    fname = fname.replace('(', '\(')
                    fname = fname.replace(')', '\)')
                    fname = fname.replace(' ', '\ ')
                    msg = 'cp ' + str(fname) + ' ' + str(img_folder)
                    os.system(msg)
                    findings[order_list[currentRow]].modify_screenshots(currentRow)
                    self.listScreenshots.addItem(new_image)
                    current_list = [str(self.listScreenshots.item(i).text()) for i in range(self.listScreenshots.count())]
                    index = len(current_list)-1
                    self.listScreenshots.item(index).setSelected(True)
                    self.listScreenshots.setCurrentRow(index)


        #MOVE ONE SCREENSHOT UP THE LIST
        def screenshotsUp():
            if self.listScreenshots.selectedItems() != []:
                currentRow = self.listScreenshots.currentRow()
                currentItem = self.listScreenshots.takeItem(currentRow)
                self.listScreenshots.insertItem(currentRow - 1, currentItem)
                self.listScreenshots.setCurrentItem(currentItem) 
                finding_item = self.listFindings.selectedIndexes()[0]
                finding_num = finding_item.row()
                findings[order_list[finding_num]].screenshots.clear()
                current_list = [str(self.listScreenshots.item(i).text()) for i in range(self.listScreenshots.count())]
                for i in current_list:
                    findings[order_list[finding_num]].screenshots.append(i)
                preview_image()
                findings[order_list[finding_num]].modify_screenshots(finding_num)


        #MOVE ONE SCREENSHOT DOWN THE LIST
        def screenshotsDown():
            if self.listScreenshots.selectedItems() != []:
                currentRow = self.listScreenshots.currentRow()
                currentItem = self.listScreenshots.takeItem(currentRow)
                self.listScreenshots.insertItem(currentRow + 1, currentItem)
                self.listScreenshots.setCurrentItem(currentItem)
                finding_item = self.listFindings.selectedIndexes()[0]
                finding_num = finding_item.row()
                findings[order_list[finding_num]].screenshots.clear()
                current_list = [str(self.listScreenshots.item(i).text()) for i in range(self.listScreenshots.count())]
                for i in current_list:
                    findings[order_list[finding_num]].screenshots.append(i)
                preview_image()
                findings[order_list[finding_num]].modify_screenshots(finding_num)


        #DELETE SCREENSHOT FROM LIST
        def delScreenshot():
            global screenshots
            if self.listScreenshots.selectedItems() != []:
                item = self.listScreenshots.takeItem(self.listScreenshots.currentRow())
                item = None
                finding_item = self.listFindings.selectedIndexes()[0]
                finding_num = finding_item.row()
                findings[order_list[finding_num]].screenshots.clear()
                current_list = [str(self.listScreenshots.item(i).text()) for i in range(self.listScreenshots.count())]
                for i in current_list:
                    findings[order_list[finding_num]].screenshots.append(i)
                findings[order_list[finding_num]].modify_screenshots(finding_num)


        #PowerPoint Presentation
        def create_pptx():
            global event
            #OPEN/CREATE PRESENTATION
            i_file = cwd + '/ERB_Template.pptx'
            prs = Presentation(i_file)
            datetime = QDateTime.currentDateTime()
            datetime = datetime.toString()

            #OUTPUT FILE NAME
            if event.event_type == 'CVPA':
                prefix = 'CVPA_'
            if event.event_type == 'PMR':
                prefix = 'PMR_'
            #SLIDE DECK MODE - LIGHT / DARK
            if event.mode == 'DARK':
                slide_deck = [9, 10, 11, 12, 13, 14, 15, 16, 17]
                c_color = [255, 255, 255]
                pptx_file = prefix + 'ERB_dark_' + datetime + '.pptx'
                pptx_file = pptx_file.replace(":", "_")
                pptx_file = pptx_file.replace(" ", "_")
            else:
                slide_deck = [0, 1, 2, 3, 4, 5, 6, 7, 8]
                c_color = [0, 0, 0]
                pptx_file = prefix + 'ERB_white_' + datetime + '.pptx'
                pptx_file = pptx_file.replace(":", "_")
                pptx_file = pptx_file.replace(" ", "_")

            #LAYOUTS
            cover_slide_layout = prs.slide_layouts[slide_deck[0]] #COVER SLIDE
            cover_slide_cui_layout = prs.slide_layouts[slide_deck[1]] #COVER SLIDE (CUI)
            title_subtitle_slide_layout = prs.slide_layouts[slide_deck[2]] #TITLE & SUBTITLE
            title_paragraph_bullets_slide_layout = prs.slide_layouts[slide_deck[3]] #TITLE & PARAGRAPH BULLETS
            title_paragraph_no_bullets_slide_layout = prs.slide_layouts[slide_deck[4]] #TITLE & PARAGRAPH NO BULLETS
            blank_slide_layout = prs.slide_layouts[slide_deck[5]] #BLANK SLIDE WITH TITLE
            picture_slide_layout = prs.slide_layouts[slide_deck[6]] #PICTURE SLIDE
            pen_process_slide_layout = prs.slide_layouts[slide_deck[7]] #PENETRATION PROCESS SLIDE
            resiliency_analysis_slide_layout = prs.slide_layouts[slide_deck[8]] #CYBER RESILIENCY ANALYSIS SLIDE

            #CLASSIFICATION MARKINGS
            long_gen_marking = 'UNCLASSIFIED//FOR OFFICIAL USE ONLY'
            short_gen_marking = '(CUI)'
            gen_color = [45, 145, 45] #green
            if event.classification == 'UNCLASSIFIED':
                color = [45, 145, 45] #green
                if event.designation == 'CUI':
                    long_marking = 'CUI'
                    short_marking = '(CUI)'
                    if event.draft == 'YES':
                        slide_marking = 'CUI//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'CUI'
                elif event.designation == 'FOUO':
                    long_marking = 'UNCLASSIFIED//FOUO'
                    short_marking = '(U//FOUO)'
                    short_gen_marking = '(U//FOUO)'
                    if event.draft == 'YES':
                        slide_marking = 'UNCLASSIFIED//FOR OFFICIAL USE ONLY//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'UNCLASSIFIED//FOR OFFICIAL USE ONLY'
                else:
                    long_marking = 'UNCLASSIFIED'
                    short_marking = '(U)'
                    long_gen_marking = 'UNCLASSIFIED'
                    short_gen_marking = '(U)'
                    if event.draft == 'YES':
                        slide_marking = 'UNCLASSIFIED//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'UNCLASSIFIED'
            if event.classification == 'SECRET':
                color = [255, 0, 0] #red
                if event.designation == 'NOFORN':
                    long_marking = 'SECRET//NOFORN'
                    short_marking = '(S//NOFORN)'
                    if event.draft == 'YES':
                        slide_marking = 'SECRET//NOFORN//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'SECRET//NOFORN'
                else:
                    long_marking = 'SECRET'
                    short_marking = '(S)'
                    if event.draft == 'YES':
                        slide_marking = 'SECRET//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'SECRET'
            if event.classification == 'TOP SECRET':
                if event.designation == 'SCI':
                    color = [186, 186, 0] #yellow
                    long_marking = 'TOP SECRET//SCI'
                    short_marking = '(TS//SCI)'
                    if event.draft == 'YES':
                        slide_marking = 'TOP SECRET//SCI//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'TOP SECRET//SCI'
                else:
                    color = [255, 140, 0] #orange
                    long_marking = 'TOP SECRET'
                    short_marking = '(TS)'
                    if event.draft == 'YES':
                        slide_marking = 'TOP SECRET//DRAFT//PRE-DECISIONAL'
                    else:
                        slide_marking = 'TOP SECRET'

            def set_classification(body, marking, colors, alignment):
                tf = body.text_frame
                p = tf.paragraphs[0]
                if alignment == 'left':
                    p.alignment = PP_ALIGN.LEFT
                elif alignment == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = marking
                font = run.font
                font.name = 'Arial'
                font.size = Pt(8)
                font.bold = True
                font.color.rgb = RGBColor(colors[0], colors[1], colors[2])

            def set_distribution(body, text1, text2):
                tf = body.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.JUSTIFY
                run = p.add_run()
                run.text = text1
                font = run.font
                font.name = 'Arial'
                font.size = Pt(8)
                font.bold = True
                run = p.add_run()
                run.text = text2
                font = run.font
                font.name = 'Arial'
                font.size = Pt(8)
                font.bold = False
            
            def set_cui(body, text):
                tf = body.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = text
                font = run.font
                font.name = 'Arial'
                font.size = Pt(8)
                font.bold = False

            def set_title(shape, text):
                title = shape.title
                title.text = text

            def set_text(body, text):
                body.text = text
            
            def set_bullet(body, text, level):
                tf = body.text_frame
                p = tf.add_paragraph()
                p.text = text
                p.level = level
                
            def set_findings(body, attrib, text, mode):
                tf = body.text_frame
                p = tf.add_paragraph()
                if mode == 2:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = attrib
                font = run.font
                font.bold = True
                run = p.add_run()
                run.text = str(text)
                font = run.font
                font.bold = False
                if mode == 0:
                    p = tf.add_paragraph()
            
            def set_info(body, text, f_bold, f_size, c_color):
                tf = body.text_frame
                p = tf.add_paragraph()
                p.text = text
                p.font.bold = f_bold
                p.font.size = Pt(f_size)
                p.font.color.rgb = RGBColor(c_color[0], c_color[1], c_color[2])
                p.alignment = PP_ALIGN.CENTER

            def _add_image(slide, placeholder_id, image_url):
                placeholder = slide.placeholders[placeholder_id]
                im = Image.open(image_url)
                #REMOVE DEFAULT IMAGE PLACEHOLDER
                image = slide.shapes[1]
                sp = image.element
                sp.getparent().remove(sp)
                #INSERT IMAGE WITH NO DEFAULT PLACEHOLDER
                im_width, im_height = im.size
                image_ratio = (im_width)/float(1.0) / (im_height)/float(1.0)
                fixed_width = 9.70
                fixed_height = 5.20
                fixed_ratio = (fixed_width)/float(1.0) / (fixed_height)/float(1.0)
                top = Inches(1.62)
                if image_ratio > fixed_ratio:
                    new_width = Inches(fixed_width)
                    new_height = Inches(fixed_width / image_ratio)
                    #ADJUST MARKINGS PLACEHOLDERS VERTICALLY
                    mid_height = (fixed_height - (fixed_width / image_ratio))/float(2.0)
                    #TOP
                    shapes.placeholders[17].top = Inches(1.24 + mid_height)
                    shapes.placeholders[17].left = Inches(2.44)
                    shapes.placeholders[17].width = Inches(2.5)
                    shapes.placeholders[17].height = Inches(0.39)
                    #BOTTOM
                    shapes.placeholders[18].top = Inches(6.81 - mid_height)
                    shapes.placeholders[18].left = Inches(2.44)
                    shapes.placeholders[18].width = Inches(2.5)
                    shapes.placeholders[18].height = Inches(0.39)
                    #ADJUST IMAGE VERTICAL POSITION
                    top = Inches(1.62 + mid_height)
                else:
                    new_width = Inches(fixed_height * image_ratio)
                    new_height = Inches(fixed_height)
                left = Inches(5) - (new_width/2)
                pic = shapes.add_picture(image_url, left, top, width=new_width, height=new_height)
                #ADJUST MARKINGS PLACEHOLDERS HORIZONTALLY
                #TOP
                shapes.placeholders[17].top = int(top - Inches(0.39))
                shapes.placeholders[17].left = int(left - Inches(0.12))
                shapes.placeholders[17].width = Inches(2.5)
                shapes.placeholders[17].height = Inches(0.39)
                #BOTTOM
                shapes.placeholders[18].top = int(top + new_height)
                shapes.placeholders[18].left = int(left + new_width - Inches(2.5) + Inches(0.12))
                shapes.placeholders[18].width = Inches(2.5)
                shapes.placeholders[18].height = Inches(0.39)
                #ADD FRAME TO SCREENSHOT
                line = pic.line
                line.color.rgb = RGBColor(0xFF, 0xDA, 0x3D)
                line.width = Pt(3)


            #COVER SLIDE
            if event.designation == 'CUI' or event.classification == 'SECRET' or event.classification == 'TOP SECRET':
                slide = prs.slides.add_slide(cover_slide_cui_layout)
            else:
                slide = prs.slides.add_slide(cover_slide_layout)
            shapes = slide.shapes
            #TITLE OF PROJECT
            if event.event_name != '':
                text = '(U) ' + event.event_name
            else:
                text = '(U) SUBTITLE GOES HERE'
            set_text(shapes.placeholders[12], text)
            if event.event_type == 'CVPA':
                set_text(shapes.placeholders[22], '(U) Emerging Results Brief (ERB)')
            else:
                set_text(shapes.placeholders[22], '(U) Cyber Resiliency Outbrief – Prevent, Mitigate, and Recover (PMR)')
            #NAME OF LEAD
            if event.lead_name != '':
                text = event.lead_name
            else:
                text = 'Name of Presenter'
            set_text(shapes.placeholders[13], text)
            #TITLE/RANK OF LEAD
            if event.lead_title != '':
                text = event.lead_title
            else:
                text = 'Rank/Title of Presenter'
            set_text(shapes.placeholders[14], text)
            #ORGANIZATION OF LEAD
            if event.lead_org != '':
                text = event.lead_org
            else:
                text = 'Organization of Presenter'
            set_text(shapes.placeholders[15], text)
            #DATE
            end_date = self.endDate.date()
            end_date = end_date.toString("dd MMM yyyy")
            if end_date != '':
                text = end_date
            else:
                text = 'DD MMM YYYY'
            set_text(shapes.placeholders[17], text)
            #CUI STATEMENT
            if event.designation == 'CUI' or event.classification == 'SECRET' or event.classification == 'TOP SECRET':
                office_symbol = event.office_symbol[0:10]
                text = 'Controlled by: ' + office_symbol + '\nCUI Category: DCRIT, Export Control\nDistribution/Dissemination Controls: D\nPOC: ' + event.lead_name + ', (575) 678-xxxx'
                set_cui(shapes.placeholders[24], text)
            #DISTRIBUTION STATEMENT
            dist_date = self.endDate.date().longMonthName(self.endDate.date().month()) + ' ' + str(self.endDate.date().year())
            text1 = '            DISTRIBUTION STATEMENT D. '
            text2 = 'Distribution authorized to the Department of Defense and U.S. DOD contractors only; administrative or operational use; export control; vulnerability information (' + dist_date + '). Other requests for this document shall be referred to Director, U.S. Army DEVCOM Analysis Center, White Sands Missile Range, NM 88002.'
            set_distribution(shapes.placeholders[23], text1, text2)
            #CLASSIFICATION - SLIDE
            set_classification(shapes.placeholders[20], slide_marking, color, 'center')
            set_classification(shapes.placeholders[21], slide_marking, color, 'center')


            if event.event_type == 'CVPA':
                #SCOPE SLIDE
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) SCOPE')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], short_marking + ' <<Add Scope description as needed>>')
                set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
                set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' ', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
                set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #SYSTEM UNDER TEST
                slide = prs.slides.add_slide(picture_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) SYSTEM UNDER TEST')
                #CLASSIFICATION - FIGURE
                set_classification(shapes.placeholders[17], long_marking, color, 'left')
                set_classification(shapes.placeholders[18], long_marking, color, 'right')
                #ADJUST MARKINGS PLACEHOLDERS HORIZONTALLY
                #TOP
                shapes.placeholders[17].top = int(Inches(1.64 - 0.39))
                shapes.placeholders[17].left = int(Inches(0.88 - 0.12))
                shapes.placeholders[17].width = Inches(2.5)
                shapes.placeholders[17].height = Inches(0.39)
                #BOTTOM
                shapes.placeholders[18].top = int(Inches(1.64 + 5.14))
                shapes.placeholders[18].left = int(Inches(0.88 + 8.27 - 2.5 + 0.12))
                shapes.placeholders[18].width = Inches(2.5)
                shapes.placeholders[18].height = Inches(0.39)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #AGENDA - EXECUTED ACTIVITIES SLIDE
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) AGENDA - EXECUTED ACTIVITIES')
                #WORK DAYS
                work_days = []
                day_1 = self.startDate.date()
                str_day_1 = day_1.toString("dd MMM yyyy")
                work_days.append(str_day_1)
                for i in range(1,5):
                    next_day = day_1.addDays(i)
                    str_next_day = next_day.toString("dd MMM yyyy")
                    work_days.append(str_next_day)
                #DAY 1
                set_text(shapes.placeholders[1], short_gen_marking + ' ' + work_days[0])
                text = '(U) In-processing, setup, network connectivity testing, discovery and enumeration scans, started collecting DOT&E metrics, and began the penetration test.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                #DAY 2
                set_bullet(shapes.placeholders[1], short_gen_marking + ' ' + work_days[1], 0)
                text = '(U) Continuation of the penetration test and DOT&E metric collection.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                #DAY 3
                set_bullet(shapes.placeholders[1], short_gen_marking + ' ' + work_days[2], 0)
                text = '(U) Continuation of the penetration test, DOT&E metric collection, and Personnel Interviews.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                #DAY 4
                set_bullet(shapes.placeholders[1], short_gen_marking + ' ' + work_days[3], 0)
                text = '(U) Continuation of the penetration test and DOT&E metric collection.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                #DAY 5
                set_bullet(shapes.placeholders[1], short_gen_marking + ' ' + work_days[4], 0)
                text = '(U) Completed the penetration test, performed the system cleanup and restoration, data consolidation, and backup. Performed the Emerging Results Brief (ERB) presentation to stakeholders.'
                set_bullet(shapes.placeholders[1], text, 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #PENETRATION TESTING ACTIVITIES
                slide = prs.slides.add_slide(pen_process_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) ACTIVITIES COMPLETION STATUS')
                #CLASSIFICATION - FIGURE
                set_classification(shapes.placeholders[17], 'CUI', gen_color, 'left')
                set_classification(shapes.placeholders[22], 'CUI', gen_color, 'right')
                #ADJUST MARKINGS PLACEHOLDERS HORIZONTALLY
                #TOP
                shapes.placeholders[17].top = Inches(0.90)
                shapes.placeholders[17].left = Inches(0.22)
                shapes.placeholders[17].width = Inches(2.5)
                shapes.placeholders[17].height = Inches(0.39)
                #BOTTOM
                shapes.placeholders[22].top = Inches(6.94)
                shapes.placeholders[22].left = Inches(7.20)
                shapes.placeholders[22].width = Inches(2.5)
                shapes.placeholders[22].height = Inches(0.39)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #PENETRATION TESTING PROCESS SLIDE
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) PENETRATION TESTING PROCESS')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], '(U) Characterization of key cyber terrain and attack vector generation')
                text = '(U) Documentation review, OSINT, site visit, staff interview, identify cyber postures, and develop attack vectors.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Discovery and Enumeration Scans'
                set_bullet(shapes.placeholders[1], text, 0)
                text = '(U) Map network, automated scanning for well-known weaknesses.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Penetration Testing'
                set_bullet(shapes.placeholders[1], text, 0)
                text = '(U) Manual probing, exploration, data pillaging, lateral movement.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Risk Analysis'
                set_bullet(shapes.placeholders[1], text, 0)
                text = '(U) Assess impact to confidentiality, integrity, and availability.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Mitigation and Risk Reduction Strategies'
                set_bullet(shapes.placeholders[1], text, 0)
                text = '(U) Develop and provide potential mitigation and risk reduction strategies to the discovered vulnerabilities.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Follow-on Testing'
                set_bullet(shapes.placeholders[1], text, 0)
                text = '(U) After mitigations are implemented, re-test to ensure the fixes are effective and do not introduce new vulnerabilities.'
                set_bullet(shapes.placeholders[1], text, 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #POSTURES SLIDE
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) POSTURES')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], '(U) Findings in this ERB constitute raw results and the technical risk analysis has not been determined.')
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) All technical findings assume some level of physical or logical access to the assets.'
                set_bullet(shapes.placeholders[1], text, 0)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Each finding will be from a specific posture. We define this postures to be as follow:'
                set_bullet(shapes.placeholders[1], text, 0)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Insider – is a person with legitimate access to the system, both logical (credentialed user) and physical or remote access.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Nearsider – physical access is provided to the target network and system, but with no credentials given.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Outsider – is a person without legitimate physical and logical access to the system under test and it is placed outside the accreditation boundary. The outsider posture is normally portrayed by an actor pivoting off a system that is legitimate connected external vectors such as SIPRNet, or Sensors.'
                set_bullet(shapes.placeholders[1], text, 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #FINDINGS TABLE SLIDE - 12 FINDINGS PER SLIDE MAX
                limit = 12
                per_slide = []
                if len(order_list) > 0:
                    reminder = len(order_list)%limit
                    complete = int(len(order_list)/limit)
                    if complete > 0:
                        for i in range(complete):
                            per_slide.append(12)
                    if reminder != 0:
                        per_slide.append(reminder)
                    for i in range(len(per_slide)):
                        slide = prs.slides.add_slide(picture_slide_layout)
                        shapes = slide.shapes
                        #REMOVE DEFAULT IMAGE PLACEHOLDER
                        image = slide.shapes[1]
                        sp = image.element
                        sp.getparent().remove(sp)
                        #CLASSIFICATION - FIGURE
                        set_classification(shapes.placeholders[17], long_marking, color, 'left')
                        set_classification(shapes.placeholders[18], long_marking, color, 'right')
                        #TITLE
                        set_title(shapes, '(U) Table of Findings')
                        offset = 2.7 - (per_slide[i]*0.1)
                        t_left = Inches(1.2)
                        t_offset = offset - 0.28
                        t_top = Inches(t_offset)
                        #ADJUST MARKINGS PLACEHOLDERS
                        #TOP
                        shapes.placeholders[17].top = int(t_top - Inches(0.05))
                        shapes.placeholders[17].left = int(t_left - Inches(0.07))
                        shapes.placeholders[17].width = Inches(2.5)
                        shapes.placeholders[17].height = Inches(0.39)
                        #BOTTOM
                        calc = offset + 0.425*(per_slide[i]+1) + 0.05
                        t_top = Inches(calc)
                        shapes.placeholders[18].top = int(t_top - Inches(0.035*(per_slide[i]+1)))
                        shapes.placeholders[18].left = int(t_left + Inches(7.5/2) + Inches(2.5/2) + Inches(0.09))
                        shapes.placeholders[18].width = Inches(2.5)
                        shapes.placeholders[18].height = Inches(0.39)
                        #TABLE
                        rows = per_slide[i]+1
                        cols = 2
                        left = Inches(1.2)
                        top = Inches(offset)
                        width = Inches(6.0)
                        height = Inches(0.8)
                        table_1 = shapes.add_table(rows, cols, left, top, width, height).table
                        #SET COLUMN WIDTHS
                        table_1.columns[0].width = Inches(0.5)
                        table_1.columns[1].width = Inches(7.0)
                        #SET COLUMN HEADINGS
                        table_1.cell(0,1).text = 'Findings'
                        #WRITE BODY CELLS
                        for j in range(per_slide[i]):
                            index = j + (i*limit)
                            table_1.cell(j+1,0).text = str(index+1)
                            table_1.cell(j+1,1).text = findings[order_list[index]].title
                            if len(table_1.cell(j+1,1).text) > 62:
                                cell = table_1.rows[j+1].cells[1]
                                paragraph = cell.text_frame.paragraphs[0]
                                paragraph.font.size = Pt(10)
                        #CLASSIFICATION
                        set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                        set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #FINDINGS - DESCRIPTION & SCREENSHOTS
                for i in order_list:
                    #FINDING SLIDE
                    slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                    shapes = slide.shapes
                    #TITLE
                    set_title(shapes, short_marking + ' ' + findings[i].title)
                    #POSTURE
                    set_findings(shapes.placeholders[1], short_marking + ' Posture: ', findings[i].posture.capitalize(), 0)
                    #AFFECTED SYSTEMS
                    set_findings(shapes.placeholders[1], short_marking + ' Affected System(s): ', findings[i].hosts, 0)
                    #ISSUES
                    set_findings(shapes.placeholders[1], short_marking + ' Issue(s): ', findings[i].issues, 1)
                    #MITIGATION
                    if findings[i].include_mitigation == 'yes':
                        set_findings(shapes.placeholders[1], short_marking + ' Mitigation: ', findings[i].mitigation, 2)
                    #CLASSIFICATION
                    set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                    set_classification(shapes.placeholders[21], slide_marking, color, 'center')
                    #SCREENSHOTS
                    if findings[i].screenshots != None:
                        for j in findings[i].screenshots:
                            ss_folder = findings[i].folder
                            i_file = ss_folder + j
                            slide = prs.slides.add_slide(picture_slide_layout)
                            shapes = slide.shapes
                            #TITLE
                            set_title(shapes, short_marking + ' ' + findings[i].title)
                            #FIGURE
                            _add_image(slide, 10, i_file)
                            #CLASSIFICATION - FIGURE
                            set_classification(shapes.placeholders[17], long_marking, color, 'left')
                            set_classification(shapes.placeholders[18], long_marking, color, 'right')
                            #CLASSIFICATION - SLIDE
                            set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                            set_classification(shapes.placeholders[21], slide_marking, color, 'center')


            if event.event_type == 'PMR':
                #CYBER RESILIENCY OBJECTIVE SLIDE
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) CYBER RESILIENCY OBJECTIVE')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], '(CUI) In support of <<SYSTEM / PROGRAM>> Evaluation and Assessments')
                text = '(CUI) Data collected during the event will inform about the cyber resilience posture.'
                set_bullet(shapes.placeholders[1], text, 1)
                set_bullet(shapes.placeholders[1], '', 0)
                text = '(U) Guided by Director, Operational Test and Evaluation (DOT&E) Memo'
                set_bullet(shapes.placeholders[1], text, 0)
                text = '(U) “Procedure For Operational Test And Evaluation Of Cybersecurity In Acquisition Programs” April 3, 2018'
                set_bullet(shapes.placeholders[1], text, 1)
                text = '(U) Prevent – The ability to protect critical mission functions from cyber threats.'
                set_bullet(shapes.placeholders[1], text, 2)
                text = '(U) Mitigate – The ability to detect and respond to cyber-attacks, and assess resilience to survive attacks and complete critical missions and tasks.'
                set_bullet(shapes.placeholders[1], text, 2)
                text = '(U) Recover – The resilience to recover from cyber-attacks and prepare mission systems”'
                set_bullet(shapes.placeholders[1], text, 2)
                text1 = '(CUI) OUR GOAL:'
                text2 = ' Provide the data and analysis on the cyber resiliency of the system to stakeholders and determine how the overall Defense Cyber Operations Team (DCOT) performed in a cyber-contested environment. Assist stakeholders to understand and have a clear picture of gaps and their strengths in DCO. Assist the program by proposing solutions and mitigations to further increase the cyber-robustness of the systems.'
                tf = shapes.placeholders[1].text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = text1
                font = run.font
                font.size = Pt(14)
                font.bold = True
                run = p.add_run()
                run.text = text2
                font = run.font
                font.size = Pt(14)
                font.bold = False
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #PMR CRITERIA - PREVENT
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) PMR CRITERIA')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], '(U) PREVENT:')
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '(U) Defined as an action that is typically proactive and implemented before the occurrence of a cyber-threat activity.', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '(U) An activity that is part of an attack vector executed by TSMO that does not meet the intended objective/effect due to a preventative mechanism in place will result in a successful prevention.', 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #PMR CRITERIA - MITIGATE
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) PMR CRITERIA')
                #BODY WITH BULLETS
                text1 = '(U) MITIGATE –'
                text2 = ' Activities are divided into two subcategories:'
                tf = shapes.placeholders[1].text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = text1
                font = run.font
                font.size = Pt(16)
                font.bold = True
                run = p.add_run()
                run.text = text2
                font = run.font
                font.size = Pt(16)
                font.bold = False
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '(U) DETECT', 1)
                set_bullet(shapes.placeholders[1], '(U) A successful detect is accomplished when a defender/operator acknowledges the technical/system detect activity and consequently creates a report (cyber incident report/helpdesk ticket), this can be triggered by an automated alert or manual review of logs. A human detect can also be accomplished by an operator/defender detecting abnormal activity on their system without the need for a network defense or monitoring tool.', 2)
                set_bullet(shapes.placeholders[1], '(U) REACT', 1)
                set_bullet(shapes.placeholders[1], '(U) A successful reaction is defined as force or action that successfully counteracts, hinders, thwarts, and/or mitigates the cyber threat action (i.e. denying the threat access to a service, port, and/or host).', 2)
                set_bullet(shapes.placeholders[1], '', 0)
                text1 = '(U) Note:'
                text2 = ' Mitigate metrics as specified in Attachment C of the DOT&E document will be based off of information extracted from logs, incident reports, observer logs, help desk tickets, and other products from the defender that will be parsed for the information requested under the measurements column of the mitigate actions. PMR observer logs will be used to fill-in any gaps left by the defender when necessary.'
                tf = shapes.placeholders[1].text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = text1
                font = run.font
                font.size = Pt(14)
                font.bold = True
                run = p.add_run()
                run.text = text2
                font = run.font
                font.size = Pt(14)
                font.bold = False
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #PMR CRITERIA - RECOVER
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) PMR CRITERIA')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], '(U) RECOVER:')
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '(U) Recover activities are those taken by operators and/or network defenders to restore mission and/or technical capabilities to continue operations after a degradation of such capabilities.', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '(U) Recover activities include, but are not limited to, re-imaging a host and/or failover to alternate sites.', 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


                #CYBER RESILIENCY ANALYSIS
                slide = prs.slides.add_slide(resiliency_analysis_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) CYBER RESILIENCY ANALYSIS')
                #CLASSIFICATION - FIGURE
                set_classification(shapes.placeholders[17], 'CUI', gen_color, 'left')
                set_classification(shapes.placeholders[22], 'CUI', gen_color, 'right')
                #ADJUST MARKINGS PLACEHOLDERS HORIZONTALLY
                #TOP
                shapes.placeholders[17].top = Inches(1.80)
                shapes.placeholders[17].left = Inches(0.22)
                shapes.placeholders[17].width = Inches(2.5)
                shapes.placeholders[17].height = Inches(0.39)
                #BOTTOM
                shapes.placeholders[22].top = Inches(5.95)
                shapes.placeholders[22].left = Inches(7.20)
                shapes.placeholders[22].width = Inches(2.5)
                shapes.placeholders[22].height = Inches(0.39)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')

                #WHITE CARDS
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) WHITECARDS')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], short_marking + ' The following Whitecards were granted to the TSMO AA Team:')
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' (1) mm/dd/yyyy @ hh:mm – <<Add Whitecard description as needed>>', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' (2) mm/dd/yyyy @ hh:mm – <<Add Whitecard description as needed>>', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' (3) mm/dd/yyyy @ hh:mm – <<Add Whitecard description as needed>>', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' (4) mm/dd/yyyy @ hh:mm – <<Add Whitecard description as needed>>', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' (5) mm/dd/yyyy @ hh:mm – <<Add Whitecard description as needed>>', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')

            #OVERALL OBSERVATIONS
            slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
            shapes = slide.shapes
            #TITLE
            if event.event_type == 'CVPA':
                set_title(shapes, '(U) OVERALL OBSERVATIONS')
            else:
                set_title(shapes, '(U) CYBER RESILIENCY TEAM OBSERVATIONS')
            #BODY WITH BULLETS
            set_text(shapes.placeholders[1], short_marking + ' System Strengths')
            set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
            set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
            set_bullet(shapes.placeholders[1], '', 0)
            set_bullet(shapes.placeholders[1], '', 0)
            set_bullet(shapes.placeholders[1], short_marking + ' System Weaknesses', 0)
            set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
            set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
            set_bullet(shapes.placeholders[1], '', 0)
            set_bullet(shapes.placeholders[1], '', 0)
            if event.event_type == 'CVPA':
                set_bullet(shapes.placeholders[1], short_marking + ' Overall Mitigations', 0)
                set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
                set_bullet(shapes.placeholders[1], short_marking + ' ', 1)
            #CLASSIFICATION - SLIDE
            set_classification(shapes.placeholders[20], slide_marking, color, 'center')
            set_classification(shapes.placeholders[21], slide_marking, color, 'center')


            if event.event_type == 'CVPA':
                #POST ASSESSMENT REPORTING
                slide = prs.slides.add_slide(title_paragraph_bullets_slide_layout)
                shapes = slide.shapes
                #TITLE
                set_title(shapes, '(U) POST ASSESSMENT REPORTING')
                #BODY WITH BULLETS
                set_text(shapes.placeholders[1], '')
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_gen_marking + ' Emerging Results Brief (ERB)', 0)
                set_bullet(shapes.placeholders[1], '(U) List of findings with minimal analysis', 1)
                set_bullet(shapes.placeholders[1], '(U) Overall assessment objective completion status', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_gen_marking + ' Assessment Results Matrix (ARM)', 0)
                set_bullet(shapes.placeholders[1], '(U) List of findings with technical risk levels', 1)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], '', 0)
                set_bullet(shapes.placeholders[1], short_gen_marking + ' Technical Report', 0)
                set_bullet(shapes.placeholders[1], '(U) Option A --> Technical Memorandum ~30 working days', 1)
                set_bullet(shapes.placeholders[1], '(U) Option B --> Published Report ~90 working days', 1)
                #CLASSIFICATION - SLIDE
                set_classification(shapes.placeholders[20], slide_marking, color, 'center')
                set_classification(shapes.placeholders[21], slide_marking, color, 'center')


            #LAST SLIDE - CONTACT INFORMATION SLIDE
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes
            #TITLE
            set_title(shapes, '(U) Contact Information')
            left = Inches(1)
            top = Inches(2)
            width = height = Inches(8)
            #BODY WITH NO BULLETS
            txBox = slide.shapes.add_textbox(left, top, width, height)
            set_text(txBox.text_frame, '')
            set_info(txBox, event.lead_name, True, 28, c_color)
            set_info(txBox, 'UNCLASSIFIED: xxxxxxxx.civ@army.mil', False, 20, c_color)
            set_info(txBox, 'SIPR: xxxxxxxx.civ@mail.smil.mil', False, 20, c_color)
            set_info(txBox, 'O: (575) xxx-xxxx', False, 20, c_color)
            set_info(txBox, 'M: (575) xxx-xxxx', False, 20, c_color)
            #CLASSIFICATION - SLIDE
            set_classification(shapes.placeholders[20], slide_marking, color, 'center')
            set_classification(shapes.placeholders[21], slide_marking, color, 'center')


            #SAVE DOCUMENT
            o_file = desktop_dir + '/' + pptx_file
            prs.save(o_file)
            #MESSAGE BOX!
            msg = QMessageBox()
            msg.setIcon(QMessageBox.Information)
            msg.setWindowTitle("PowerPoint")
            txt_msg = "PowerPoint successfully created at " + o_file + "\nThank you... Good bye!"
            msg.setText(txt_msg)
            x = msg.exec_()
            #EXIT
            app.quit()


        #SWITCH BETWEEN PAGES USING STACKED WIDGETS / CREATE PMR SLIDES
        def go_page2():
            global event, folder_list
            if self.cvpaButton.isChecked() == True:
                event_type = 'CVPA'
            else:
                event_type = 'PMR'
            lead_name = self.leadName.text() #lead name
            lead_title = self.leadTitle.text() #lead title
            lead_org = self.leadOrg.text() #lead org
            office_symbol = self.office_comboBox.currentText()
            event_name = self.eventName.text() #event name
            start_date = self.startDate.text() #start date
            end_date = self.endDate.text() #end date
            if self.unclassButton.isChecked() == True:
                classification = 'UNCLASSIFIED' #classification - unclassified
                if self.cuiButton.isChecked() == True:
                    designation = 'CUI' #designation - cui
                elif self.fouoButton.isChecked() == True:
                    designation = 'FOUO' #designation - fouo
                else:
                    designation = 'NONE' #designation - none
            if self.secretButton.isChecked() == True:
                classification = 'SECRET' #classification - secret
                if self.noforncheckBox.isChecked() == True:
                    designation = 'NOFORN' #designation - noforn
                else:
                    designation = 'NONE' #designation - none
            if self.topsecretButton.isChecked() == True:
                classification = 'TOP SECRET' #classification - top secret
                if self.scicheckBox.isChecked() == True:
                    designation = 'SCI' #designation - sci
                else:
                    designation = 'NONE' #designation - none
            if self.draftcheckBox.isChecked() == True:
                draft = 'YES' #draft - yes
            else:
            	draft = 'NO' #draft - no
            if self.lightButton.isChecked() == True:
                mode = 'LIGHT' #mode - light
            if self.darkButton.isChecked() == True:
                mode = 'DARK' #mode - dark
            event = Event(lead_name, lead_title, lead_org, office_symbol, event_name, event_type, start_date, end_date, classification, designation, draft, mode)
            event.write_file()
            if self.cvpaButton.isChecked() == True:
                if '.zip' in data_source:
                    aux_data_folder = data_folder.replace(' ', '\ ')
                    msg = 'python3 dradis_parser.py ' + aux_data_folder
                    msg = msg.replace('(', '\(')
                    msg = msg.replace(')', '\)')
                    os.system(msg)
                if 'fric_export_' in data_source:
                    msg = 'python3 fric_parser.py ' + data_folder
                    os.system(msg)
                if 'Create your own ERB' in data_source:
                    msg = 'python3 empty_erb.py'
                    os.system(msg)
                current_selection = self.listFolders.currentItem().text()
                if folder_list != current_selection:
                    get_findings()
                self.stackedWidget0.setCurrentIndex(1)
            if self.pmrButton.isChecked() == True:
                create_pptx()
            

        def go_page1():
            global folder_list
            folder_list = self.listFolders.currentItem().text()
            self.stackedWidget0.setCurrentIndex(0)


        self.cvpaButton.toggled.connect(event_type)
        self.pmrButton.toggled.connect(event_type)
        self.unclassButton.toggled.connect(class_changed)
        self.secretButton.toggled.connect(class_changed)
        self.topsecretButton.toggled.connect(class_changed)
        self.listFolders.itemSelectionChanged.connect(on_selection_changed)
        self.listFindings.itemSelectionChanged.connect(set_fields)
        self.listScreenshots.itemSelectionChanged.connect(preview_image)
        self.nextButton.clicked.connect(go_page2)
        self.gobackButton.clicked.connect(go_page1)
        self.findingAddButton.clicked.connect(addFinding)
        self.findingUpButton.clicked.connect(findingsUp)
        self.findingDownButton.clicked.connect(findingsDown)
        self.findingDeleteButton.clicked.connect(delFinding)
        self.updateDescButton.clicked.connect(update_finding)
        self.screenshotAddButton.clicked.connect(addScreenshot)
        self.screenshotUpButton.clicked.connect(screenshotsUp)
        self.screenshotDownButton.clicked.connect(screenshotsDown)
        self.screenshotDeleteButton.clicked.connect(delScreenshot)
        self.pptxButton.clicked.connect(create_pptx)
        self.quitButton.clicked.connect(QApplication.instance().quit)


    def retranslateUi(self, MainWindow):
        _translate = QtCore.QCoreApplication.translate
        MainWindow.setWindowTitle(_translate("MainWindow", "Emerging Results Brief (ERB) Generator"))
        self.startDate.setDisplayFormat(_translate("MainWindow", "MM/dd/yyyy"))
        self.leadName_label.setText(_translate("MainWindow", "Name:"))
        self.leadName.setPlaceholderText(_translate("MainWindow", "Your name here!"))
        self.leadTitle.setText(_translate("MainWindow", "Computer Scientist"))
        self.eventName.setPlaceholderText(_translate("MainWindow", "Title of the Event"))
        self.leadTitle_label.setText(_translate("MainWindow", "Rank/Title:"))
        self.leadOrg_label.setText(_translate("MainWindow", "Organization:"))
        self.endDate.setDisplayFormat(_translate("MainWindow", "MM/dd/yyyy"))
        self.startDate_label.setText(_translate("MainWindow", "Start Date:"))
        self.endDate_label.setText(_translate("MainWindow", "End Date:"))
        self.dradisfricData_label.setText(_translate("MainWindow", "DRADIS/FRIC DATA"))
        self.eventName_label.setText(_translate("MainWindow", "Name:"))
        self.leadOrg.setText(_translate("MainWindow", "DEVCOM Analysis Center"))
        self.event_label.setText(_translate("MainWindow", "EVENT"))
        self.teamLead_label.setText(_translate("MainWindow", "TEAM LEAD"))
        self.availFilesFolders_label.setText(_translate("MainWindow", "Available Files / Folders:"))
        self.nextButton.setText(_translate("MainWindow", "Next"))
        self.unclassButton.setText(_translate("MainWindow", "UNCLASSIFIED"))
        self.secretButton.setText(_translate("MainWindow", "SECRET"))
        self.topsecretButton.setText(_translate("MainWindow", "TOP SECRET"))
        self.noforncheckBox.setText(_translate("MainWindow", "NOFORN"))
        self.scicheckBox.setText(_translate("MainWindow", "SCI"))
        self.fouoButton.setText(_translate("MainWindow", "FOUO (obsolete)"))
        self.cuiButton.setText(_translate("MainWindow", "CUI"))
        self.mode_label.setText(_translate("MainWindow", "Slides\n"
"Background"))
        self.lightButton.setText(_translate("MainWindow", "LIGHT"))
        self.darkButton.setText(_translate("MainWindow", "DARK"))
        self.classification_label.setText(_translate("MainWindow", "Classification"))
        self.draftcheckBox.setText(_translate("MainWindow", "DRAFT//PRE-DECISIONAL"))
        self.office_label.setText(_translate("MainWindow", "Office Symbol"))
        self.pmrButton.setText(_translate("MainWindow", "PMR"))
        self.cvpaButton.setText(_translate("MainWindow", "CVPA / CVI"))
        self.eventType_label.setText(_translate("MainWindow", "Type:"))
        self.quitButton.setText(_translate("MainWindow", "Quit"))
        self.findings_label.setText(_translate("MainWindow", "Findings:"))
        self.updateDescButton.setText(_translate("MainWindow", "Update Finding"))
        self.screenshotPreview_label.setText(_translate("MainWindow", "Screenshot Preview:"))
        self.imagePreview_label.setText(_translate("MainWindow", "NO PREVIEW"))
        self.findingDeleteButton.setText(_translate("MainWindow", "DEL"))
        self.findingUpButton.setText(_translate("MainWindow", "UP"))
        self.findingDownButton.setText(_translate("MainWindow", "DOWN"))
        self.screenshots_label.setText(_translate("MainWindow", "Screenshots:"))
        self.pptxButton.setText(_translate("MainWindow", "PPTX"))
        self.issues_label.setText(_translate("MainWindow", "Issues:"))
        self.gobackButton.setText(_translate("MainWindow", "Go Back"))
        self.screenshotDeleteButton.setText(_translate("MainWindow", "DEL"))
        self.screenshotDownButton.setText(_translate("MainWindow", "DOWN"))
        self.screenshotUpButton.setText(_translate("MainWindow", "UP"))
        self.findingName_label.setText(_translate("MainWindow", "Finding Name:"))
        self.findingHosts_label.setText(_translate("MainWindow", "Affected Hosts:"))
        self.posture_label.setText(_translate("MainWindow", "Posture:"))
        self.insiderButton.setText(_translate("MainWindow", "Insider"))
        self.nearsiderButton.setText(_translate("MainWindow", "Nearsider"))
        self.outsiderButton.setText(_translate("MainWindow", "Outsider"))
        self.mitigation_label.setText(_translate("MainWindow", "Mitigation:"))
        self.mitigationcheckBox.setText(_translate("MainWindow", "Include Mitigation"))
        self.findingAddButton.setText(_translate("MainWindow", "ADD"))
        self.screenshotAddButton.setText(_translate("MainWindow", "ADD"))


if __name__ == "__main__":
    import sys
    def exitHandler():
        try:
            #REMOVE INACTIVE FINDINGS IN XML
            tree = ET.parse(current_erb)
            root = tree.getroot()
            find_to_del = []
            for finding in root.findall('finding'):
                uid = finding.get('uid')
                active = int(finding.find('active').text)
                if active == 0:
                    root.remove(finding)
                    folder = finding.find('folder')
                    find_to_del.append(folder.text)
            tree.write(current_erb)
            #UPDATE AND SORT FINDINGS BY RANK
            tree = ET.parse(current_erb)
            root = tree.getroot()
            findings_dict = {}
            for finding in root.findall('finding'):
                uid = finding.get('uid')
                rank = int(finding.find('rank').text)
                findings_dict[uid] = rank
            sorted_dict = sorted(findings_dict.items(), key=lambda kv:kv[1])
            new_dict = {}
            for i in range(len(sorted_dict)):
                new_dict[sorted_dict[i][0]] = i
            for finding in root.findall('finding'):
                uid = finding.get('uid')
                rank = finding.find('rank')
                rank.text = str(new_dict[uid])
            tree.write(current_erb)
            #UPDATE UIDs
            tree = ET.parse(current_erb)
            root = tree.getroot()
            sorted_dict = sorted(new_dict.items(), key=lambda kv:int(kv[0]))
            new_dict = {}
            for i in range(len(sorted_dict)):
                new_dict[sorted_dict[i][0]] = sorted_dict[i][1]
            for finding in root.findall('finding'):
                uid = finding.get('uid')
                index = list(new_dict.keys()).index(uid)
                finding.set('uid', str(index))
            tree.write(current_erb)
            #DELETE UNUSED SCREENSHOTS
            for i in find_to_del:
                if os.path.exists(i):
                    os.system('rm -r ' + i)
        except:
            pass
    app = QtWidgets.QApplication(sys.argv)
    app.aboutToQuit.connect(exitHandler)
    MainWindow = QtWidgets.QMainWindow()
    ui = Ui_MainWindow()
    ui.setupUi(MainWindow)
    MainWindow.show()
    sys.exit(app.exec_())
